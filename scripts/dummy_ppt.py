from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches
from random import choice
import string

# Create a new PowerPoint presentation
prs = Presentation()

# Add a slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add some random text to the slide
for i in range(5):
    left = Inches(i*0.5)
    top = Inches(0.5)
    width = Inches(1)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    words = "".join(choice(string.ascii_letters) for i in range(10))
    tf.text = words
    
# Add a WordArt to the slide
shapes = slide.shapes
word_art_shape = shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
word_art_shape.text_frame.clear()
word_art_shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
word_art_shape.word_art.text = "WordArt"
word_art_shape.word_art.font.size = Pt(24)

# Save the presentation
prs.save("dummy.pptx")
